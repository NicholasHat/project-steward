"""PPTX via python-pptx: per-slide text (titles kept separately as
structure) and document core properties."""

from __future__ import annotations

from pathlib import Path

from truth_engine.parse.registry import register
from truth_engine.parse.types import ParsedDocument

PARSER_NAME = "python-pptx"
PARSER_VERSION = "1"


@register("pptx")
def parse_pptx(path: Path) -> ParsedDocument:
    from pptx import Presentation  # lazy: pipeline extra

    presentation = Presentation(str(path))

    slide_titles: list[str | None] = []
    slide_texts: list[str] = []
    for slide_num, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title else None
        slide_titles.append(title)

        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if parts:
            slide_texts.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

    core = presentation.core_properties
    embedded_metadata = {
        k: v
        for k, v in {
            "author": core.author,
            "title": core.title,
            "created": core.created.isoformat() if core.created else None,
            "modified": core.modified.isoformat() if core.modified else None,
        }.items()
        if v
    }

    return ParsedDocument(
        raw_text="\n\n".join(slide_texts) if slide_texts else None,
        structure={"slide_count": len(presentation.slides), "slide_titles": slide_titles},
        embedded_metadata=embedded_metadata or None,
        parser_name=PARSER_NAME,
        parser_version=PARSER_VERSION,
    )
