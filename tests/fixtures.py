"""Tiny sample-document builders for parser tests.

PDFs are built by hand (byte offsets computed from the actual buffer, not
memorized/hardcoded) rather than shelling out to a platform-specific tool
(e.g. macOS's `cupsfilter`) or adding a PDF-writing dependency the app itself
never needs — this keeps the fixture portable and dependency-free. The other
formats use the same libraries the corresponding handler reads with, via
their write APIs.
"""

from __future__ import annotations

from pathlib import Path


def make_pdf(path: Path, text: str) -> None:
    """A minimal single-page PDF containing `text`, readable by pdfplumber."""
    objects: list[bytes] = []

    def obj(n: int, body: bytes) -> bytes:
        return f"{n} 0 obj\n".encode() + body + b"\nendobj\n"

    content_stream = f"BT /F1 12 Tf 20 100 Td ({text}) Tj ET".encode()
    objects.append(obj(1, b"<< /Type /Catalog /Pages 2 0 R >>"))
    objects.append(obj(2, b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>"))
    objects.append(
        obj(
            3,
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 300] "
            b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        )
    )
    objects.append(
        obj(
            4,
            f"<< /Length {len(content_stream)} >>\nstream\n".encode()
            + content_stream
            + b"\nendstream",
        )
    )
    objects.append(obj(5, b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"))

    buf = bytearray(b"%PDF-1.4\n")
    offsets = [0]  # object 0 is the conventional free-list head
    for o in objects:
        offsets.append(len(buf))
        buf += o

    xref_offset = len(buf)
    count = len(objects) + 1
    xref = [f"xref\n0 {count}\n".encode(), b"0000000000 65535 f \n"]
    xref += [f"{off:010d} 00000 n \n".encode() for off in offsets[1:]]
    buf += b"".join(xref)
    buf += f"trailer\n<< /Size {count} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF".encode()

    path.write_bytes(bytes(buf))


def make_docx(path: Path, heading: str, paragraph: str) -> None:
    import docx

    document = docx.Document()
    document.add_heading(heading, level=1)
    document.add_paragraph(paragraph)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "name"
    table.cell(0, 1).text = "value"
    table.cell(1, 0).text = "widget"
    table.cell(1, 1).text = "42"
    document.save(str(path))


def make_xlsx(path: Path, sheet_name: str, rows: list[list[object]]) -> None:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = sheet_name
    for row in rows:
        sheet.append(row)
    workbook.save(str(path))


def make_pptx(path: Path, title: str, body: str) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    presentation.save(str(path))


def make_png(path: Path, size: tuple[int, int] = (16, 16)) -> None:
    from PIL import Image

    Image.new("RGB", size, color=(200, 100, 50)).save(path, format="PNG")
